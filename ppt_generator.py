import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def extract_lesson_sections(lesson_content):
    """Extract sections from the lesson plan content"""
    sections = {}
    
    # Extract title
    title_match = re.search(r'# Lesson Plan: (.+)', lesson_content)
    if title_match:
        sections['title'] = title_match.group(1).strip()
    else:
        sections['title'] = "Lesson Plan"
    
    # Extract grade, topic, strategy
    overview_info = {}
    overview_matches = re.findall(r'\*\*(.*?):\*\* (.*?)(?=\n|\*\*|$)', lesson_content)
    for key, value in overview_matches:
        overview_info[key.lower()] = value.strip()
    
    sections['grade_level'] = overview_info.get('grade level', '')
    sections['topic'] = overview_info.get('topic', '')
    sections['strategy'] = overview_info.get('teaching strategy', '')
    
    # Extract objectives
    objectives_match = re.search(r'## Objectives\n(.*?)(?=\n##|\Z)', lesson_content, re.DOTALL)
    if objectives_match:
        objectives_text = objectives_match.group(1).strip()
        # Extract bulleted points
        bullets = re.findall(r'• (.*?)(?=\n•|\n\n|\Z)', objectives_text, re.DOTALL)
        if not bullets:  # Try without bullet points
            bullets = [line.strip() for line in objectives_text.split('\n') if line.strip()]
        sections['objectives'] = bullets
    else:
        sections['objectives'] = []
    
    # Extract introduction
    intro_match = re.search(r'## Introduction\n(.*?)(?=\n##|\Z)', lesson_content, re.DOTALL)
    if intro_match:
        intro_text = intro_match.group(1).strip()
        # Extract bulleted points
        bullets = re.findall(r'• (.*?)(?=\n•|\n\n|\Z)', intro_text, re.DOTALL)
        if not bullets:  # Try without bullet points
            bullets = [line.strip() for line in intro_text.split('\n') if line.strip()]
        sections['introduction'] = bullets
    else:
        sections['introduction'] = []
    
    # Extract main activities
    activities_match = re.search(r'## Main Activities\n(.*?)(?=\n##|\Z)', lesson_content, re.DOTALL)
    if activities_match:
        activities_text = activities_match.group(1).strip()
        # Extract bulleted points
        bullets = re.findall(r'• (.*?)(?=\n•|\n\n|\Z)', activities_text, re.DOTALL)
        if not bullets:  # Try without bullet points
            bullets = [line.strip() for line in activities_text.split('\n') if line.strip()]
        sections['activities'] = bullets
    else:
        sections['activities'] = []
    
    # Extract assessment
    assessment_match = re.search(r'## Assessment\n(.*?)(?=\n##|\Z)', lesson_content, re.DOTALL)
    if assessment_match:
        assessment_text = assessment_match.group(1).strip()
        # Extract bulleted points
        bullets = re.findall(r'• (.*?)(?=\n•|\n\n|\Z)', assessment_text, re.DOTALL)
        if not bullets:  # Try without bullet points
            bullets = [line.strip() for line in assessment_text.split('\n') if line.strip()]
        sections['assessment'] = bullets
    else:
        sections['assessment'] = []
    
    # Extract conclusion
    conclusion_match = re.search(r'## Conclusion\n(.*?)(?=\n##|\Z)', lesson_content, re.DOTALL)
    if conclusion_match:
        conclusion_text = conclusion_match.group(1).strip()
        # Extract bulleted points
        bullets = re.findall(r'• (.*?)(?=\n•|\n\n|\Z)', conclusion_text, re.DOTALL)
        if not bullets:  # Try without bullet points
            bullets = [line.strip() for line in conclusion_text.split('\n') if line.strip()]
        sections['conclusion'] = bullets
    else:
        sections['conclusion'] = []
    
    return sections

def add_title_slide(prs, title, grade_level, topic, strategy):
    """Add a title slide to the presentation"""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Set subtitle with grade, topic, and strategy
    subtitle = slide.placeholders[1]
    subtitle.text = f"Grade: {grade_level}\nTopic: {topic}\nStrategy: {strategy}"
    
    # Format title and subtitle
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(44)
            run.font.bold = True
            run.font.color.rgb = RGBColor(44, 77, 121)  # Dark blue
    
    for paragraph in subtitle.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(28)
            run.font.color.rgb = RGBColor(89, 89, 89)  # Dark gray
    
    return slide

def add_section_slide(prs, title, content_list):
    """Add a slide with a title and bulleted content"""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Format title
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(40)
            run.font.bold = True
            run.font.color.rgb = RGBColor(44, 77, 121)  # Dark blue
    
    # Add content bullets
    content = slide.placeholders[1]
    text_frame = content.text_frame
    text_frame.clear()  # Clear existing text
    
    # Add each bullet point
    for i, item in enumerate(content_list):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = item
        p.level = 0  # Top level bullet
        
        # Format bullet point text
        for run in p.runs:
            run.font.size = Pt(24)
            run.font.color.rgb = RGBColor(0, 0, 0)  # Black text
    
    return slide

def create_presentation(grade_level, topic, teaching_strategy, lesson_content, output_path):
    """Create a PowerPoint presentation from the lesson plan content"""
    # Extract sections from lesson content
    sections = extract_lesson_sections(lesson_content)
    
    # Create presentation
    prs = Presentation()
    
    # Set slide dimensions (16:9 aspect ratio)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Add title slide
    add_title_slide(
        prs, 
        sections.get('title', f"Lesson Plan: {topic}"),
        grade_level,
        topic,
        teaching_strategy.replace('_', ' ').title()
    )
    
    # Add objectives slide
    add_section_slide(prs, "Objectives", sections.get('objectives', []))
    
    # Add introduction slide
    add_section_slide(prs, "Introduction", sections.get('introduction', []))
    
    # Add main activities slide
    add_section_slide(prs, "Main Activities", sections.get('activities', []))
    
    # Add assessment slide
    add_section_slide(prs, "Assessment", sections.get('assessment', []))
    
    # Add conclusion slide
    add_section_slide(prs, "Conclusion", sections.get('conclusion', []))
    
    # Add summary slide
    summary = [
        f"Grade Level: {grade_level}",
        f"Topic: {topic}",
        f"Teaching Strategy: {teaching_strategy.replace('_', ' ').title()}",
        "Key Points:",
        "• Review objectives and key concepts",
        "• Implement assessment strategies",
        "• Follow up with extension activities"
    ]
    add_section_slide(prs, "Summary", summary)
    
    # Save presentation
    prs.save(output_path)
    
    return output_path
